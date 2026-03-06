# ============================================
# UNIVERSAL MULTILINGUAL AI CHATBOT
# PART 1 — CORE ENGINE
# ============================================

import sys
import os
import streamlit as st
import numpy as np
import faiss
import io
import json
import time
import requests
import base64

from bs4 import BeautifulSoup
from sentence_transformers import SentenceTransformer
from langdetect import detect
from gtts import gTTS
from faster_whisper import WhisperModel

import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation

from google.oauth2 import service_account
from googleapiclient.discovery import build


# ============================================
# LANGUAGE DETECTION
# ============================================

def detect_language(text):

    try:

        lang = detect(text)

        mapping = {
            "en": "English",
            "hi": "Hindi",
            "te": "Telugu",
            "fr": "French"
        }

        return mapping.get(lang, "English")

    except:

        return "English"


# ============================================
# TEXT TO SPEECH
# ============================================

def speak_response(text, language):

    lang_codes = {
        "English": "en",
        "Hindi": "hi",
        "Telugu": "te",
        "French": "fr"
    }

    tts = gTTS(text=text, lang=lang_codes.get(language, "en"))

    filename = "response.mp3"

    tts.save(filename)

    audio_file = open(filename, "rb").read()

    b64 = base64.b64encode(audio_file).decode()

    audio_html = f"""
    <audio autoplay>
    <source src="data:audio/mp3;base64,{b64}" type="audio/mp3">
    </audio>
    """

    st.markdown(audio_html, unsafe_allow_html=True)


# ============================================
# WHISPER SPEECH RECOGNITION
# ============================================

@st.cache_resource
def load_whisper():

    model = WhisperModel("base", compute_type="int8")

    return model


whisper_model = load_whisper()


def transcribe_audio(file):

    segments, info = whisper_model.transcribe(file)

    text = ""

    for seg in segments:

        text += seg.text

    return text


# ============================================
# AI PROVIDERS
# ============================================

class AIProvider:

    def get_response(self, prompt):

        raise NotImplementedError


class OpenAIProvider(AIProvider):

    def __init__(self, api_key):

        import openai

        self.client = openai.OpenAI(api_key=api_key)

    def get_response(self, prompt):

        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1200
        )

        return response.choices[0].message.content


class GeminiProvider(AIProvider):

    def __init__(self, api_key):

        import google.generativeai as genai

        genai.configure(api_key=api_key)

        self.model = genai.GenerativeModel("gemini-2.5-flash")

    def get_response(self, prompt):

        response = self.model.generate_content(prompt)

        return response.text


class ClaudeProvider(AIProvider):

    def __init__(self, api_key):

        import anthropic

        self.client = anthropic.Anthropic(api_key=api_key)

    def get_response(self, prompt):

        response = self.client.messages.create(
            model="claude-3-sonnet-20240229",
            max_tokens=1200,
            messages=[{"role": "user", "content": prompt}]
        )

        return response.content[0].text


# ============================================
# MULTILINGUAL RAG CHATBOT
# ============================================

class UniversalChatbot:

    def __init__(self, ai_provider, credentials_json=None):

        self.ai_provider = ai_provider

        # Multilingual embedding model
        self.embedding_model = SentenceTransformer(
            "intfloat/multilingual-e5-base"
        )

        self.documents = []
        self.index = None
        self.embeddings = None

        self.SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

        self.service = None

        if credentials_json:

            self._setup_service_account(credentials_json)


    # ============================================
    # GOOGLE DRIVE CONNECTION
    # ============================================

    def _setup_service_account(self, credentials_json):

        credentials = service_account.Credentials.from_service_account_info(
            credentials_json,
            scopes=self.SCOPES
        )

        self.service = build("drive", "v3", credentials=credentials)
# ============================================
# DOCUMENT EXTRACTION
# ============================================

    def extract_text_from_pdf(self, content):

        pdf_file = io.BytesIO(content)

        reader = PyPDF2.PdfReader(pdf_file)

        text = ""

        for page in reader.pages[:15]:

            text += page.extract_text() + "\n"

        return text


    def extract_text_from_word(self, content):

        doc_file = io.BytesIO(content)

        doc = Document(doc_file)

        text = ""

        for p in doc.paragraphs:

            text += p.text + "\n"

        return text


    def extract_text_from_excel(self, content):

        excel_file = io.BytesIO(content)

        df = pd.read_excel(excel_file)

        return df.to_string()


    def extract_text_from_ppt(self, content):

        ppt_file = io.BytesIO(content)

        prs = Presentation(ppt_file)

        text = ""

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

        return text


# ============================================
# WEBSITE SCRAPER
# ============================================

    def scrape_website(self, url):

        headers = {

            "User-Agent": "Mozilla/5.0"

        }

        response = requests.get(url, headers=headers)

        soup = BeautifulSoup(response.text, "html.parser")

        for tag in soup(["script", "style", "nav", "footer"]):

            tag.decompose()

        text = soup.get_text()

        return text[:12000]


# ============================================
# CHUNKING
# ============================================

    def chunk_text(self, text, size=800, overlap=150):

        chunks = []

        start = 0

        while start < len(text):

            end = start + size

            chunk = text[start:end]

            chunks.append(chunk)

            start = end - overlap

        return chunks


# ============================================
# BUILD VECTOR DATABASE
# ============================================

    def build_index(self):

        all_chunks = [

            f"passage: {doc['content']}"

            for doc in self.documents

        ]

        embeddings = self.embedding_model.encode(all_chunks)

        dimension = embeddings.shape[1]

        self.index = faiss.IndexFlatL2(dimension)

        self.index.add(np.array(embeddings).astype("float32"))

        self.embeddings = embeddings


# ============================================
# SEARCH
# ============================================

    def search(self, query, k=4):

        query_embedding = self.embedding_model.encode(

            [f"query: {query}"]

        )

        scores, ids = self.index.search(

            np.array(query_embedding).astype("float32"),
            k
        )

        results = []

        for i in ids[0]:

            results.append(self.documents[i]["content"])

        return results


# ============================================
# RESPONSE GENERATION
# ============================================

    def get_response(self, question):

        user_lang = detect_language(question)

        context_docs = self.search(question)

        context = "\n\n".join(context_docs)

        prompt = f"""
Answer the question using the context below.

Respond ONLY in {user_lang}

Context:
{context}

Question:
{question}

Answer clearly:
"""

        response = self.ai_provider.get_response(prompt)

        return response


# ============================================
# STREAMLIT UI
# ============================================

def main():

    st.set_page_config(page_title="Universal Multilingual AI", page_icon="🤖")

    st.title("🌍 Universal Multilingual AI Chatbot")


# ============================================
# AI PROVIDER SELECTION
# ============================================

    st.sidebar.header("AI Provider")

    provider = st.sidebar.selectbox(

        "Choose AI",

        ["Gemini", "OpenAI", "Claude"]

    )

    api_key = st.sidebar.text_input("API Key", type="password")

    if provider == "Gemini":

        ai = GeminiProvider(api_key)

    elif provider == "OpenAI":

        ai = OpenAIProvider(api_key)

    else:

        ai = ClaudeProvider(api_key)


# ============================================
# INITIALIZE BOT
# ============================================

    if "chatbot" not in st.session_state:

        st.session_state.chatbot = UniversalChatbot(ai)


# ============================================
# FILE UPLOAD
# ============================================

    uploaded_files = st.sidebar.file_uploader(

        "Upload Knowledge Files",

        accept_multiple_files=True

    )

    if uploaded_files:

        bot = st.session_state.chatbot

        for file in uploaded_files:

            content = file.read()

            name = file.name

            if name.endswith(".pdf"):

                text = bot.extract_text_from_pdf(content)

            elif name.endswith(".docx"):

                text = bot.extract_text_from_word(content)

            elif name.endswith(".xlsx"):

                text = bot.extract_text_from_excel(content)

            elif name.endswith(".pptx"):

                text = bot.extract_text_from_ppt(content)

            else:

                text = content.decode()

            chunks = bot.chunk_text(text)

            for c in chunks:

                bot.documents.append({

                    "source": name,

                    "content": c

                })

        bot.build_index()

        st.success("Knowledge loaded")


# ============================================
# VOICE INPUT
# ============================================

    audio_file = st.file_uploader(

        "🎤 Upload voice question",

        type=["wav", "mp3"]

    )

    voice_prompt = None

    if audio_file:

        with open("temp_audio.wav", "wb") as f:

            f.write(audio_file.read())

        voice_prompt = transcribe_audio("temp_audio.wav")

        st.write("Speech detected:", voice_prompt)


# ============================================
# CHAT INPUT
# ============================================

    text_prompt = st.chat_input("Ask anything")

    prompt = text_prompt if text_prompt else voice_prompt


# ============================================
# CHAT PROCESS
# ============================================

    if prompt:

        bot = st.session_state.chatbot

        with st.spinner("Thinking..."):

            answer = bot.get_response(prompt)

        st.markdown(answer)

        lang = detect_language(prompt)

        speak_response(answer, lang)


if __name__ == "__main__":

    main()
